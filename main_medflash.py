# CÓDIGO FINAL DE MED-FLASH AI v2.0 (Versión Blindada Anti-Errores)
import streamlit as st
import time
import json
import random 
import yaml
from yaml.loader import SafeLoader

try:
    # --- Importaciones Críticas ---
    from PIL import Image
    import fitz  # PyMuPDF
    from pptx import Presentation
    import pandas as pd
    import google.generativeai as genai
    import plotly.graph_objects as go
    import firebase_admin
    from firebase_admin import credentials, firestore
    import streamlit_authenticator as stauth
    import bcrypt
    from streamlit_authenticator.utilities.hasher import Hasher 
except ImportError as e:
    st.error("Error crítico de dependencias.")
    st.code(f"Error: {e}")
    st.stop()

# --- CONFIGURACIÓN DE PÁGINA ---
st.set_page_config(
    page_title="Med-Flash AI v2.0",
    page_icon="🧬",
    layout="wide",
    initial_sidebar_state="collapsed", 
)

# --- INICIALIZACIÓN DE MEMORIA OFFLINE (FALLBACK) ---
if 'offline_db' not in st.session_state:
    st.session_state.offline_db = {
        'users': {},
        'decks': {}
    }

# --- VÍNCULOS VISUALES DINÁMICOS ---
SYSTEM_VISUALS = {
    "Cardiovascular": {"icon": "❤️", "color": "#FF5757"},
    "Respiratorio": {"icon": "🫁", "color": "#46B9C7"},
    "Nervioso Central": {"icon": "🧠", "color": "#A67CEF"},
    "Nervioso Periférico": {"icon": "⚡", "color": "#FFD700"},
    "Digestivo": {"icon": "🍔", "color": "#FFB347"},
    "Renal (Urinario)": {"icon": "💧", "color": "#5C94FF"},
    "Musculoesquelético": {"icon": "💪", "color": "#90EE90"},
    "Endocrino": {"icon": "🧬", "color": "#FF69B4"},
    "Hematológico": {"icon": "🩸", "color": "#DC143C"},
    "Inmunológico": {"icon": "🛡️", "color": "#1E90FF"},
    "Reproductivo": {"icon": "🤰", "color": "#F5A6C1"},
    "Metabolismo": {"icon": "🔥", "color": "#FF8C00"},
    "Enzimas/Proteínas": {"icon": "🧩", "color": "#32CD32"},
    "Genética/ADN": {"icon": "🧬", "color": "#8A2BE2"},
    "Biología Celular": {"icon": "🦠", "color": "#20B2AA"},
    "Farmacocinética": {"icon": "📈", "color": "#FFD700"},
    "Farmacodinámica": {"icon": "🎯", "color": "#FF4500"},
    "Antibióticos": {"icon": "💊", "color": "#00CED1"},
    "General": {"icon": "📚", "color": "#E0E0E0"},
    "Otro": {"icon": "❓", "color": "#4A4A4A"},
    "Seleccionar Sistema": {"icon": "🩺", "color": "#F5A6C1"},
}

MATERIAS = [
    "Seleccionar Materia", "Anatomía", "Fisiología", "Patología", "Semiología", 
    "Bioquímica", "Genética", "Biología Celular", 
    "Farmacología", "Microbiología", 
    "Pediatría", "Neurología", "Cardiología", "Medicina Interna"
]

SISTEMAS_CUERPO = [
    "Cardiovascular", "Respiratorio", "Nervioso Central", "Nervioso Periférico", 
    "Digestivo", "Renal (Urinario)", "Musculoesquelético", "Endocrino", 
    "Hematológico", "Inmunológico", "Reproductivo", "General"
]

TOPICOS_POR_MATERIA = {
    "Bioquímica": ["Metabolismo", "Enzimas/Proteínas", "Genética/ADN", "General"],
    "Genética": ["Genética/ADN", "Biología Celular", "General"],
    "Biología Celular": ["Biología Celular", "Genética/ADN", "Metabolismo"],
    "Farmacología": ["Farmacocinética", "Farmacodinámica", "Antibióticos"] + SISTEMAS_CUERPO,
    "Microbiología": ["Antibióticos", "Inmunológico", "General"],
    "DEFAULT": SISTEMAS_CUERPO
}

# --- ESTILOS CSS (CORRECCIONES DE CONTRASTE Y BRILLO) ---
st.markdown("""
<style>
    /* 1. FONDO / UI: Acuarela Artística */
    .stApp {
        background-color: #fdfbf7; /* Base papel crema muy suave */
        background-image: 
            radial-gradient(at 0% 0%, hsla(340,82%,76%,0.3) 0px, transparent 50%),
            radial-gradient(at 100% 0%, hsla(210,29%,24%,0.1) 0px, transparent 50%),
            radial-gradient(at 100% 100%, hsla(340,82%,76%,0.2) 0px, transparent 50%),
            radial-gradient(at 0% 100%, hsla(210,29%,24%,0.1) 0px, transparent 50%);
        background-attachment: fixed;
    }

    /* FIX BARRA LATERAL: Forzar fondo claro y texto oscuro legible */
    [data-testid="stSidebar"] {
        background-color: #f7f9fc; /* Gris azulado muy pálido */
        border-right: 1px solid rgba(0,0,0,0.05);
    }
    [data-testid="stSidebar"] * {
        color: #4A5568 !important; /* Gris pizarra para contraste perfecto */
    }
    [data-testid="stSidebar"] .stButton button {
        border: 1px solid #cbd5e0;
        color: #4A5568 !important;
    }
    [data-testid="stSidebar"] .stButton button:hover {
        border-color: #F5A6C1;
        background-color: white;
    }

    /* TEXTO GENERAL: "Bajar el brillo" (De negro puro a Gris Frío) */
    h1, h2, h3, p, label, .stMarkdown, .stRadio label {
        color: #4A5568 !important; /* Mucho más suave para la vista */
        font-family: 'Helvetica Neue', sans-serif;
    }
    
    /* Títulos principales un poco más oscuros para jerarquía */
    h1, h2 {
        color: #2D3748 !important;
    }

    /* 2. TARJETA DE PREGUNTA: Marco Dorado Iridiscente */
    .flashcard {
        background: rgba(255, 255, 255, 0.9);
        backdrop-filter: blur(12px);
        border-radius: 20px;
        padding: 40px;
        margin: 30px 0;
        
        /* El Borde Dorado Holográfico */
        border: 3px solid transparent;
        background-image: linear-gradient(white, white), 
                          linear-gradient(135deg, #C5A059, #F2E6C2, #C5A059); /* Oro más mate/elegante */
        background-origin: border-box;
        background-clip: content-box, border-box;
        box-shadow: 0 8px 20px rgba(197, 160, 89, 0.2); 
        
        color: #2D3748;
        font-size: 1.2rem;
        font-weight: 500;
    }

    /* 3. CONTENIDO CENTRAL (FEEDBACK): Hiperrealismo 8K */
    .feedback-container {
        margin-top: 25px;
        border-radius: 12px;
        overflow: hidden;
        box-shadow: 0 20px 40px -10px rgba(0, 0, 0, 0.8);
        border: 1px solid rgba(255, 255, 255, 0.1);
    }

    .feedback-correct {
        background: linear-gradient(145deg, #142615, #0a110a);
        color: #86efac; /* Verde suave, menos neón */
        padding: 20px;
        font-weight: bold;
        border-bottom: 1px solid rgba(134, 239, 172, 0.1);
    }

    .feedback-incorrect {
        background: linear-gradient(145deg, #2b1111, #140808);
        color: #fca5a5; /* Rojo suave, menos neón */
        padding: 20px;
        font-weight: bold;
        border-bottom: 1px solid rgba(252, 165, 165, 0.1);
    }

    .feedback-explanation {
        background: #1a202c; /* Gris muy oscuro (Gunmetal) en vez de negro total */
        color: #edf2f7; /* Blanco hueso para texto */
        padding: 30px;
        font-family: 'Segoe UI', sans-serif;
        line-height: 1.7;
    }
    
    .feedback-explanation table {
        width: 100%;
        border-collapse: collapse;
        margin: 15px 0;
        background: #2d3748;
        border-radius: 8px;
        overflow: hidden;
    }
    .feedback-explanation th {
        background: #4a5568;
        color: #F6E05E; 
        padding: 10px;
    }
    .feedback-explanation td {
        padding: 10px;
        border-bottom: 1px solid #4a5568;
    }

    /* Doodle Container */
    .doodle-container {
        width: 100%; height: 150px; 
        background-color: white; 
        border-radius: 16px; 
        display: flex; flex-direction: column; 
        align-items: center; justify-content: center; 
        margin-bottom: 20px; padding: 10px; 
        border: 2px solid var(--system-color, #d69e2e);
        box-shadow: 0 2px 10px rgba(0,0,0,0.03);
    }
    .doodle-container .system-icon { font-size: 4rem; line-height: 1; }
    .doodle-container .system-text { color: #4A5568 !important; font-weight: bold; font-size: 0.85rem; }
    
    /* Botones Principales */
    .stButton > button {
        border-radius: 25px;
        background-image: linear-gradient(to right, #Eeb4c9 0%, #e6a6be 100%); /* Rosa más mate */
        color: white !important;
        border: none;
        font-weight: 600;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        transition: all 0.2s;
    }
    .stButton > button:hover {
        transform: translateY(-1px);
        box-shadow: 0 6px 10px rgba(0,0,0,0.15);
    }
    
</style>
""", unsafe_allow_html=True)

# --- Funciones Auxiliares ---
def extraer_texto_pdf(file_stream):
    try:
        doc = fitz.open(stream=file_stream.read(), filetype="pdf")
        texto = ""
        for page in doc: texto += page.get_text()
        doc.close()
        return texto
    except Exception as e: return f"Error PDF: {e}"

def extraer_texto_pptx(file_stream):
    try:
        prs = Presentation(file_stream)
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): texto += shape.text + "\n"
        return texto
    except Exception as e: return f"Error PPTX: {e}"

# --- Estado de Sesión ---
if 'page' not in st.session_state: st.session_state.page = "Cargar Contenido"
if 'extracted_content' not in st.session_state: st.session_state.extracted_content = None
if 'current_exam' not in st.session_state: st.session_state.current_exam = None
if 'current_question_index' not in st.session_state: st.session_state.current_question_index = 0
if 'user_answer' not in st.session_state: st.session_state.user_answer = None
if 'show_explanation' not in st.session_state: st.session_state.show_explanation = False
if 'exam_results' not in st.session_state: st.session_state.exam_results = []
if "authentication_status" not in st.session_state: st.session_state.authentication_status = None
if "user_level" not in st.session_state: st.session_state.user_level = "Nivel 1 (Novato)"
if "materia_actual" not in st.session_state: st.session_state.materia_actual = MATERIAS[0]
if "sistema_actual" not in st.session_state: st.session_state.sistema_actual = "General"

def restart_exam():
    st.session_state.current_exam = None
    st.session_state.current_question_index = 0
    st.session_state.user_answer = None
    st.session_state.show_explanation = False
    st.session_state.exam_results = []

def go_to_next_question():
    st.session_state.current_question_index += 1
    st.session_state.user_answer = None
    st.session_state.show_explanation = False

# --- API & Database ---
@st.cache_resource
def init_firebase():
    # 1. Verificación de existencia
    if "FIREBASE_SERVICE_ACCOUNT" not in st.secrets:
        return None
    try:
        secret_value = st.secrets["FIREBASE_SERVICE_ACCOUNT"].strip()
        try:
            cred_json = json.loads(secret_value)
        except json.JSONDecodeError as e:
            st.error(f"⚠️ Error de formato en 'Secrets': {e}")
            return None
        cred = credentials.Certificate(cred_json)
        if not firebase_admin._apps: 
            firebase_admin.initialize_app(cred)
        return firestore.client()
    except Exception as e:
        st.error(f"⚠️ Error conectando a Firebase: {e}")
        return None

db = init_firebase()

api_key_disponible = "GOOGLE_API_KEY" in st.secrets and st.secrets["GOOGLE_API_KEY"]
gemini_model = None
if api_key_disponible:
    try:
        genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])
        gemini_model = genai.GenerativeModel(model_name="gemini-2.5-flash-preview-09-2025")
    except Exception as e:
        pass

# --- CAPA DE DATOS HÍBRIDA ---
def get_all_users_credentials():
    """Carga usuarios de DB, memoria offline y puente de memoria."""
    try:
        test_hash = Hasher(['123']).generate()[0] 
    except Exception:
        test_hash = "$2b$12$y.X.1.1.1.1.1.1.1.1.1.u.1.1.1.1.1.1.1.1.1.1.1.1.1.1.1"

    base_credentials = {
        'usernames': {
            'drdavid': {
                'name': 'Dr. David',
                'email': 'david@medflash.ai',
                'password': test_hash,
                'progreso': {},
                'logged_in': False
            }
        }
    }

    # 1. Cargar desde Firebase (Si existe)
    if db:
        try:
            users_ref = db.collection('usuarios')
            docs = users_ref.stream()
            for doc in docs:
                data = doc.to_dict()
                if 'password' in data:
                    base_credentials['usernames'][doc.id] = data
        except Exception as e:
            print(f"Error DB: {e}")

    # 2. Sobrescribir/Añadir con memoria local (Latency Fix)
    if 'offline_db' in st.session_state:
        for u, data in st.session_state.offline_db['users'].items():
            base_credentials['usernames'][u] = data

    return base_credentials

def register_new_user(name, email, username, password):
    """Registra en DB y en memoria local SIMULTÁNEAMENTE."""
    if not name or not email or not username or not password:
        return "Por favor completa todos los campos."

    hashed_pw = Hasher([password]).generate()[0]
    user_data = {
        'name': name, 
        'email': email, 
        'password': hashed_pw, 
        'progreso': {} 
    }

    # Guardar SIEMPRE en memoria local primero (Latency Fix)
    if 'offline_db' not in st.session_state:
        st.session_state.offline_db = {'users': {}, 'decks': {}}
    st.session_state.offline_db['users'][username] = user_data

    # Intento Online (Si hay DB)
    if db:
        try:
            doc_ref = db.collection('usuarios').document(username)
            doc_ref.set(user_data)
        except Exception as e:
            return f"Guardado localmente, pero error en nube: {str(e)}"

    return "success"

def get_user_progress(username, materia):
    # 1. Buscar Offline
    if 'offline_db' in st.session_state and username in st.session_state.offline_db['users']:
        progreso = st.session_state.offline_db['users'][username].get('progreso', {})
        if materia in progreso: return progreso[materia]['level'], progreso[materia]['xp']

    # 2. Buscar Online
    if not db: return "Nivel 1 (Novato)", 0
    try:
        doc = db.collection('usuarios').document(username).get()
        if doc.exists: 
            data = doc.to_dict()
            progreso = data.get('progreso', {})
            if materia in progreso: return progreso[materia]['level'], progreso[materia]['xp']
    except: pass
    return "Nivel 1 (Novato)", 0

def update_user_level(username, materia, passed):
    levels = ["Nivel 1 (Novato)", "Nivel 2 (Estudiante)", "Nivel 3 (Interno)", "Nivel 4 (Residente)", "Nivel 5 (Especialista)"]
    
    def calc_next_level(current_prog):
        if materia in current_prog:
            lvl = current_prog[materia]['level']; xp = current_prog[materia]['xp']
        else:
            lvl = "Nivel 1 (Novato)"; xp = 0
        
        if passed:
            xp += 10
            idx = levels.index(lvl) if lvl in levels else 0
            if idx < 4: 
                return levels[idx+1], xp, f"¡Subiste de nivel en {materia}! Ahora eres: {levels[idx+1]} 🌟"
        return lvl, xp, ""

    # 1. Actualizar Offline
    if 'offline_db' in st.session_state and username in st.session_state.offline_db['users']:
        user = st.session_state.offline_db['users'][username]
        nl, nx, m = calc_next_level(user.get('progreso', {}))
        if 'progreso' not in user: user['progreso'] = {}
        user['progreso'][materia] = {'level': nl, 'xp': nx}

    # 2. Actualizar Online
    if not db: 
        return nl, m if 'nl' in locals() else ""
        
    try:
        doc_ref = db.collection('usuarios').document(username)
        data = doc_ref.get().to_dict() or {}
        progreso = data.get('progreso', {})
        nl, nx, m = calc_next_level(progreso)
        progreso[materia] = {'level': nl, 'xp': nx}
        doc_ref.update({'progreso': progreso})
        return nl, m
    except: return None, None

def get_user_decks(username):
    decks = {}
    if 'offline_db' in st.session_state and username in st.session_state.offline_db['decks']:
        decks.update(st.session_state.offline_db['decks'][username])
    
    if db:
        try:
            stream = db.collection('usuarios').document(username).collection('mazos').stream()
            for d in stream:
                decks[d.id] = d.to_dict()
        except: pass
    return decks

def save_user_deck(username, name, content, mat, sis):
    deck_data = {'preguntas': content, 'materia': mat, 'sistema': sis, 'creado': str(time.time())}
    
    if 'offline_db' not in st.session_state: st.session_state.offline_db = {'users':{}, 'decks':{}}
    if 'decks' not in st.session_state.offline_db: st.session_state.offline_db['decks'] = {}
    if username not in st.session_state.offline_db['decks']: st.session_state.offline_db['decks'][username] = {}
    st.session_state.offline_db['decks'][username][name] = deck_data

    if db:
        try:
            db.collection('usuarios').document(username).collection('mazos').document(name).set({
                'preguntas': content, 'materia': mat, 'sistema': sis, 'creado': firestore.SERVER_TIMESTAMP
            })
        except: return False 
    
    return True

def delete_user_deck(username, name):
    success = False
    if 'offline_db' in st.session_state and username in st.session_state.offline_db['decks']:
        if name in st.session_state.offline_db['decks'][username]:
            del st.session_state.offline_db['decks'][username][name]
            success = True

    if db:
        try:
            db.collection('usuarios').document(username).collection('mazos').document(name).delete()
            success = True
        except: pass
    
    return success

# --- AUTHENTICATOR SETUP ---
credentials_data = get_all_users_credentials()
config = {
    'credentials': credentials_data,
    'cookie': {'expiry_days': 30, 'key': 'medflash_key_v2', 'name': 'medflash_cookie_v2'},
    'preauthorized': {'emails': []}
}
authenticator = stauth.Authenticate(
    config['credentials'], config['cookie']['name'], config['cookie']['key'], 
    config['cookie']['expiry_days'], config['preauthorized']['emails']
)

# --- MAIN APP ---
if st.session_state["authentication_status"] is None:
    st.markdown("<h1 style='text-align: center; color: #4A5568;'>Med-Flash AI v2.0 🧬</h1>", unsafe_allow_html=True)
    
    if not db:
        st.warning("⚠️ Modo Offline Activado: Datos temporales.")
    
    tab1, tab2 = st.tabs(["Login", "Registro"])
    with tab1: 
        authenticator.login('main')
        with st.expander("Verificar Usuarios Disponibles"):
            st.caption("Usuarios cargados en memoria (Login permitido):")
            st.code(list(credentials_data['usernames'].keys()))

    with tab2:
        with st.form("reg"):
            st.write("Crea tu cuenta:")
            u = st.text_input("Usuario (Este usarás para entrar)", placeholder="ej. drdavid")
            p = st.text_input("Contraseña", type="password")
            n = st.text_input("Nombre Completo", placeholder="ej. Dr. David")
            e = st.text_input("Email")
            
            if st.form_submit_button("Registrar"):
                res = register_new_user(n, e, u, p)
                if res == "success": 
                    st.success(f"¡Cuenta creada para '{u}'! Ve a la pestaña Login.")
                    time.sleep(1)
                    st.rerun()
                else: st.error(res)

elif st.session_state["authentication_status"]:
    username = st.session_state.get("username")
    name = st.session_state.get("name")
    
    materia_display = st.session_state.materia_actual
    if materia_display == "Seleccionar Materia":
        nivel_actual = "Selecciona Materia"
    else:
        l, x = get_user_progress(username, materia_display)
        nivel_actual = l
        st.session_state.user_level = nivel_actual 

    if st.session_state.get("last_login") != username:
        st.session_state.flashcard_library = get_user_decks(username)
        st.session_state.last_login = username

    current_system = st.session_state.sistema_actual
    visuals = SYSTEM_VISUALS.get(current_system, SYSTEM_VISUALS["Otro"])
    
    with st.sidebar:
        st.markdown(f"### Dr. {name}")
        if not db:
            st.caption("⚠️ MODO OFFLINE")
            
        if materia_display != "Seleccionar Materia":
            st.caption(f"Nivel en {materia_display}:")
            st.info(f"{nivel_actual}")
        else:
            st.caption("Selecciona una materia para ver tu nivel.")
            
        authenticator.logout('Salir', 'sidebar')
        st.markdown("---")
        st.markdown(f"""
        <div class="doodle-container" style="--system-color: {visuals['color']};">
            <span class="system-icon">{visuals['icon']}</span>
            <span class="system-text">{st.session_state.materia_actual}</span>
            <span class="system-text">{current_system}</span>
        </div>
        """, unsafe_allow_html=True)
        st.markdown("---")
        if st.button("1. Cargar Contenido", use_container_width=True): st.session_state.page = "Cargar Contenido"
        if st.button("2. Verificación IA", use_container_width=True): st.session_state.page = "Verificación IA"
        if st.button("3. Generar Examen", use_container_width=True): st.session_state.page = "Generar Examen"
        if st.button("4. Estudiar", use_container_width=True): st.session_state.page = "Mi Progreso"

    # --- PÁGINA 1: CARGAR ---
    if st.session_state.page == "Cargar Contenido":
        st.header("1. Contexto Clínico 📚")
        c1, c2 = st.columns(2)
        with c1:
            mat = st.selectbox("Materia:", MATERIAS)
            st.session_state.materia_actual = mat
            if mat != materia_display: st.rerun()
        with c2:
            if mat in TOPICOS_POR_MATERIA: ops = TOPICOS_POR_MATERIA[mat]
            elif mat == "Seleccionar Materia": ops = ["Selecciona Materia Primero"]
            else: ops = TOPICOS_POR_MATERIA["DEFAULT"]
            sis = st.selectbox("Tema/Sistema:", ops)
            st.session_state.sistema_actual = sis
            
        st.divider()
        f = st.file_uploader("Sube PDF/PPTX/TXT", ["pdf", "pptx", "txt"])
        if st.button("Procesar Archivo", type="primary"):
            if f:
                with st.spinner("Leyendo..."):
                    if f.type == "application/pdf": t = extraer_texto_pdf(f)
                    elif "presentation" in f.type: t = extraer_texto_pptx(f)
                    else: t = f.read().decode("utf-8")
                    st.session_state.extracted_content = t
                    st.success("Texto extraído. Continúa a 'Verificación IA'.")

    # --- PÁGINA 2: VERIFICACIÓN IA ---
    elif st.session_state.page == "Verificación IA":
        st.header("2. Verificación Médica con IA 🔬")
        if not st.session_state.extracted_content: st.warning("Carga un archivo primero."); st.stop()
        
        st.info(f"Analizando contenido de **{st.session_state.materia_actual} / {st.session_state.sistema_actual}**")
        st.text_area("Contenido:", st.session_state.extracted_content[:2000]+"...", height=200)
        
        if st.button("🔬 Analizar Precisión Científica", type="primary"):
            if not gemini_model:
                st.error("❌ Error: No se detectó la API Key de Google en los secrets.")
                st.stop()
                
            prompt = [
                f"Rol: Profesor de medicina experto en {st.session_state.materia_actual}.",
                f"Contexto: {st.session_state.materia_actual} - {st.session_state.sistema_actual}.",
                f"Texto a revisar:\n{st.session_state.extracted_content[:15000]}",
                "Tarea: Evalúa la precisión científica y claridad.",
                "Usa formato Markdown:",
                "- 🟢 Puntos Clave Correctos.",
                "- 🟡 Ambigüedades o puntos a mejorar.",
                "- 🔴 Errores potenciales o falta de contexto.",
                "Provee un resumen ejecutivo para el estudiante."
            ]
            
            with st.spinner("La IA está auditando el contenido..."):
                try:
                    response = gemini_model.generate_content(prompt)
                    st.markdown("### Informe de Auditoría IA")
                    st.markdown(response.text)
                except Exception as e:
                    st.error(f"Error en análisis: {e}")

    # --- PÁGINA 3: GENERAR EXAMEN ---
    elif st.session_state.page == "Generar Examen":
        st.header("3. Generar Flashcards Visuales 🧠")
        if not st.session_state.extracted_content: st.warning("Carga un archivo primero."); st.stop()
        
        d_name = st.text_input("Nombre del Mazo (ej. Parcial Bioquímica)")
        num = st.slider("Preguntas", 1, 10, 5)
        
        if st.button("🚀 Crear con Feedback Visual", type="primary"):
            if not gemini_model:
                st.error("❌ Error Crítico: No se detectó la API Key.")
                st.stop()

            if not d_name: st.error("Pon un nombre al mazo."); st.stop()
            restart_exam()
            
            prompt = [
                f"Eres profesor experto en {st.session_state.materia_actual} y diseñador instruccional médico.",
                f"Tema: {st.session_state.sistema_actual}. Nivel Estudiante: {st.session_state.user_level} (En {st.session_state.materia_actual}).",
                f"Texto base:\n{st.session_state.extracted_content[:10000]}...",
                f"Crea {num} preguntas de opción múltiple ADAPTADAS A ESTE NIVEL.",
                "IMPORTANTE - FEEDBACK VISUAL:",
                "En el campo 'explicacion', NO uses texto plano.",
                "Usa MARKDOWN para crear:",
                "- Tablas comparativas.",
                "- Listas con emojis (🦠, 💊, ⚡).",
                "- Diagramas de flujo de texto (A -> B -> C).",
                "- Esquemas anatómicos simples ( [Órgano] === [Tejido] ).",
                "Formato JSON array estricto:",
                """[{"pregunta": "...", "opciones": {"A": "...", "B": "...", "C": "...", "D": "..."}, "respuesta_correcta": "A", "explicacion": "Markdown rico aquí..."}]"""
            ]
            
            with st.spinner("Generando explicaciones gráficas..."):
                try:
                    res = gemini_model.generate_content(prompt)
                    txt = res.text.replace('```json', '').replace('```', '')
                    data = json.loads(txt[txt.find('['):txt.rfind(']')+1])
                    
                    if save_user_deck(username, d_name, data, st.session_state.materia_actual, st.session_state.sistema_actual):
                        if not isinstance(st.session_state.get('flashcard_library'), dict):
                             st.session_state.flashcard_library = {}
                        st.session_state.flashcard_library[d_name] = data
                        st.success("Mazo creado. Vamos a estudiar."); st.balloons()
                except Exception as e: st.error(f"Error IA: {e}")

    # --- PÁGINA 4: PROGRESO (VERSIÓN BLINDADA) ---
    elif st.session_state.page == "Mi Progreso":
        st.header("4. Biblioteca de Estudio 🏆")
        
        # 1. Recuperación y saneamiento
        raw_decks = st.session_state.get("flashcard_library", {})
        
        # Si está corrupto (no es dict), recargar desde DB
        if not isinstance(raw_decks, dict):
            raw_decks = get_user_decks(username)
            st.session_state.flashcard_library = raw_decks
        
        decks = raw_decks
        
        if not decks:
            st.info("No tienes mazos guardados aún.")
        else:
            # 2. Generación Segura de Opciones (Elemento por Elemento)
            opts = []
            deck_map = {}
            
            try:
                for k, v in decks.items():
                    # Verificamos que CADA mazo sea un diccionario válido
                    if isinstance(v, dict):
                        label = f"{k} [{v.get('materia','General')}]"
                        opts.append(label)
                        deck_map[label] = k
                
                if not opts:
                    st.warning("Tus mazos parecen vacíos o corruptos.")
                    if st.button("Limpiar Memoria"):
                        st.session_state.flashcard_library = {}
                        st.rerun()
                    st.stop()

                sel = st.selectbox("Selecciona Mazo", opts)
                real_name = deck_map[sel]
                
                c1, c2 = st.columns([1, 4])
                if c1.button("Estudiar"):
                    st.session_state.current_exam = decks[real_name]
                    st.session_state.current_exam['name'] = real_name
                    st.session_state.page = "Estudiar"
                    st.rerun()
                if c1.button("Borrar"):
                     delete_user_deck(username, real_name)
                     del st.session_state.flashcard_library[real_name]
                     st.rerun()
                     
            except Exception as e:
                st.error(f"Error de lectura: {e}")
                if st.button("Reparar Biblioteca"):
                    st.session_state.flashcard_library = get_user_decks(username)
                    st.rerun()

    # --- PÁGINA 5: ESTUDIO ---
    elif st.session_state.page == "Estudiar":
        exam = st.session_state.current_exam.get('preguntas', [])
        materia_examen = st.session_state.current_exam.get('materia', 'General')
        idx = st.session_state.current_question_index
        
        if st.button("⬅ Volver"): st.session_state.page = "Mi Progreso"; restart_exam(); st.rerun()
        if idx < len(exam):
            q = exam[idx]
            st.markdown(f"### Pregunta {idx+1}/{len(exam)}")
            
            # TARJETA DE PREGUNTA CON BORDE DORADO IRIDISCENTE
            st.markdown(f'<div class="flashcard"><h5>{q["pregunta"]}</h5></div>', unsafe_allow_html=True)
            
            ops = list(q['opciones'].values())
            sel = st.radio("Tu respuesta:", ops, key=f"q{idx}", disabled=st.session_state.show_explanation)
            if st.button("Responder") and sel:
                st.session_state.show_explanation = True
                cor_ltr = q['respuesta_correcta']
                cor_txt = q['opciones'][cor_ltr]
                is_ok = (sel == cor_txt)
                if len(st.session_state.exam_results) <= idx:
                    st.session_state.exam_results.append({'ok': is_ok, 'sel': sel, 'cor': cor_txt})
                st.rerun()

            if st.session_state.show_explanation:
                # CONTENEDOR HIPERREALISTA (FEEDBACK)
                st.markdown('<div class="feedback-container">', unsafe_allow_html=True)
                
                res = st.session_state.exam_results[idx]
                if res['ok']: 
                    st.markdown('<div class="feedback-correct">✅ ¡EXCELENTE! RESPUESTA CORRECTA</div>', unsafe_allow_html=True)
                else: 
                    st.markdown(f'<div class="feedback-incorrect">❌ INCORRECTO. RESPUESTA REAL: {res["cor"]}</div>', unsafe_allow_html=True)
                
                st.markdown(f'<div class="feedback-explanation">{q["explicacion"]}</div>', unsafe_allow_html=True)
                
                st.markdown('</div>', unsafe_allow_html=True) # Fin container
                
                st.write("") # Espacio
                if st.button("Siguiente ➡"): go_to_next_question(); st.rerun()
        else:
            st.balloons()
            score = sum(1 for r in st.session_state.exam_results if r['ok'])
            final = (score / len(exam)) * 100
            st.metric("Resultado Final", f"{final:.0f}%")
            
            nl, msg = update_user_level(username, materia_examen, final >= 80)
            if msg: st.success(msg)
            
            if materia_examen == st.session_state.materia_actual:
                st.session_state.user_level = nl if nl else st.session_state.user_level

elif st.session_state["authentication_status"] is False:
    st.error("Credenciales inválidas")
